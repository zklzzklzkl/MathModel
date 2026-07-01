#!/usr/bin/env python
"""Ingest local math-modeling knowledge files into a traceable SQLite ledger.

The script is intentionally useful without network access or heavyweight ML
dependencies. SQLite is the required baseline. ChromaDB and sentence-transformer
embeddings are optional adapters for later richer retrieval.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import math
import re
import sqlite3
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_CONFIG = REPO_ROOT / "knowledge" / "libraries.json"
DEFAULT_DB = REPO_ROOT / "knowledge" / ".local" / "rag.sqlite3"
DEFAULT_CHROMA_DIR = REPO_ROOT / "knowledge" / ".local" / "chroma"

SUPPORTED_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".csv",
    ".tsv",
    ".docx",
    ".xlsx",
    ".pptx",
    ".json",
    ".yaml",
    ".yml",
    ".tex",
    ".typ",
    ".py",
    ".r",
    ".m",
    ".jl",
    ".js",
    ".ts",
    ".sql",
    ".pdf",
}

SKIP_DIR_NAMES = {
    ".git",
    ".local",
    "__pycache__",
    "node_modules",
    ".venv",
    "dist",
    "build",
}


class IngestSkip(Exception):
    """Expected skip, such as empty files or unavailable PDF parser."""


@dataclass(frozen=True)
class Library:
    id: str
    name: str
    description: str
    stage: str
    default_license: str
    default_tags: tuple[str, ...]


@dataclass
class Chunk:
    chunk_id: str
    library_id: str
    source_path: str
    source_relpath: str
    content_hash: str
    chunk_hash: str
    chunk_index: int
    text: str
    title: str
    license: str
    year: str | None
    contest: str | None
    problem_id: str | None
    tags: list[str]
    stage: str
    recommended_use: str
    risk_warning: str


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def load_libraries(config_path: Path = DEFAULT_CONFIG) -> dict[str, Library]:
    data = json.loads(config_path.read_text(encoding="utf-8"))
    libraries: dict[str, Library] = {}
    for item in data["libraries"]:
        libraries[item["id"]] = Library(
            id=item["id"],
            name=item.get("name", item["id"]),
            description=item.get("description", ""),
            stage=item.get("stage", data.get("default_stage", "capability")),
            default_license=item.get("default_license", "local-use-only"),
            default_tags=tuple(item.get("default_tags", [])),
        )
    return libraries


def init_db(db_path: Path) -> sqlite3.Connection:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS chunks (
            chunk_id TEXT PRIMARY KEY,
            library_id TEXT NOT NULL,
            source_path TEXT NOT NULL,
            source_relpath TEXT NOT NULL,
            content_hash TEXT NOT NULL,
            chunk_hash TEXT NOT NULL,
            chunk_index INTEGER NOT NULL,
            text TEXT NOT NULL,
            title TEXT,
            license TEXT NOT NULL,
            year TEXT,
            contest TEXT,
            problem_id TEXT,
            tags_json TEXT NOT NULL,
            stage TEXT NOT NULL,
            recommended_use TEXT NOT NULL,
            risk_warning TEXT NOT NULL,
            ingested_at TEXT NOT NULL,
            UNIQUE(library_id, source_relpath, chunk_index, chunk_hash)
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ingest_errors (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source_path TEXT NOT NULL,
            library_id TEXT,
            reason TEXT NOT NULL,
            detail TEXT,
            created_at TEXT NOT NULL
        )
        """
    )
    try:
        conn.execute(
            """
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts
            USING fts5(chunk_id UNINDEXED, title, text, tags)
            """
        )
    except sqlite3.OperationalError:
        # FTS5 is bundled with most Python builds, but the lexical scorer in
        # rag_query.py remains the compatibility baseline if it is unavailable.
        pass
    return conn


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Ingest local RAG knowledge sources.")
    parser.add_argument("--source", action="append", type=Path, help="File or directory to ingest. Can be repeated.")
    parser.add_argument("--config", type=Path, default=DEFAULT_CONFIG)
    parser.add_argument("--db", type=Path, default=DEFAULT_DB)
    parser.add_argument("--library", help="Force all sources into this library id when it cannot be inferred.")
    parser.add_argument("--max-chars", type=int, default=1400)
    parser.add_argument("--min-chars", type=int, default=40)
    parser.add_argument("--vector-store", choices=["auto", "none", "chroma"], default="auto")
    parser.add_argument("--chroma-dir", type=Path, default=DEFAULT_CHROMA_DIR)
    parser.add_argument(
        "--embedding-mode",
        choices=["auto-local", "hashing", "sentence-transformer"],
        default="auto-local",
        help="Embedding mode for optional Chroma. auto-local tries locally cached BGE models, then falls back to hashing.",
    )
    parser.add_argument(
        "--embedding-model",
        default="BAAI/bge-m3",
        help="Preferred sentence-transformer model for optional Chroma indexing.",
    )
    parser.add_argument("--json", action="store_true", help="Print machine-readable summary.")
    return parser.parse_args()


def configure_stdout() -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")


def iter_source_files(paths: Iterable[Path]) -> Iterable[Path]:
    for source in paths:
        source = source.resolve()
        if source.is_file():
            if source.suffix.lower() in SUPPORTED_EXTENSIONS:
                yield source
            continue
        if not source.exists():
            continue
        for path in source.rglob("*"):
            if not path.is_file():
                continue
            if any(part in SKIP_DIR_NAMES for part in path.parts):
                continue
            if path.suffix.lower() in SUPPORTED_EXTENSIONS:
                yield path.resolve()


def decode_bytes(raw: bytes, path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "cp1252"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise IngestSkip(f"unreadable_text_encoding:{path.suffix.lower()}")


def read_pdf_text(path: Path) -> str:
    try:
        import fitz  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on optional package
        raise IngestSkip("pdf_support_requires_pymupdf") from exc
    doc = fitz.open(path)
    try:
        return "\n\n".join(page.get_text("text") for page in doc)
    finally:
        doc.close()


def read_csv_text(path: Path) -> str:
    text = decode_bytes(path.read_bytes(), path)
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    try:
        reader = csv.reader(text.splitlines(), delimiter=delimiter)
        rows = []
        for idx, row in enumerate(reader):
            if idx >= 200:
                rows.append(["... truncated after 200 rows ..."])
                break
            rows.append(row)
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
    except csv.Error:
        return text


def read_docx_text(path: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on optional package
        raise IngestSkip("docx_support_requires_python_docx") from exc
    doc = Document(path)
    parts: list[str] = []
    parts.extend(paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on optional package
        raise IngestSkip("xlsx_support_requires_openpyxl") from exc
    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets[:12]:
            parts.append(f"# Sheet: {sheet.title}")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > 200:
                    parts.append("... truncated after 200 rows ...")
                    break
                cells = ["" if cell is None else str(cell).strip() for cell in row]
                if any(cells):
                    parts.append(" | ".join(cells))
    finally:
        workbook.close()
    return "\n".join(parts)


def read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on optional package
        raise IngestSkip("pptx_support_requires_python_pptx") from exc
    presentation = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"# Slide {index}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def read_source_text(path: Path) -> str:
    if path.stat().st_size == 0:
        raise IngestSkip("empty_file")
    if path.suffix.lower() == ".pdf":
        text = read_pdf_text(path)
    elif path.suffix.lower() in {".csv", ".tsv"}:
        text = read_csv_text(path)
    elif path.suffix.lower() == ".docx":
        text = read_docx_text(path)
    elif path.suffix.lower() == ".xlsx":
        text = read_xlsx_text(path)
    elif path.suffix.lower() == ".pptx":
        text = read_pptx_text(path)
    else:
        text = decode_bytes(path.read_bytes(), path)
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        raise IngestSkip("empty_text")
    return text


def parse_scalar(value: str) -> Any:
    value = value.strip()
    if value.startswith("[") and value.endswith("]"):
        inner = value[1:-1].strip()
        if not inner:
            return []
        return [item.strip().strip("'\"") for item in inner.split(",") if item.strip()]
    if "," in value and not re.search(r"https?://", value):
        return [item.strip().strip("'\"") for item in value.split(",") if item.strip()]
    return value.strip("'\"")


def parse_front_matter(text: str) -> tuple[dict[str, Any], str]:
    if not text.startswith("---\n"):
        return {}, text
    end = text.find("\n---", 4)
    if end == -1:
        return {}, text
    raw_meta = text[4:end].strip()
    body = text[end + 4 :].lstrip("\n")
    meta: dict[str, Any] = {}
    for line in raw_meta.splitlines():
        if not line.strip() or line.lstrip().startswith("#") or ":" not in line:
            continue
        key, value = line.split(":", 1)
        meta[key.strip()] = parse_scalar(value)
    return meta, body


def infer_library_id(path: Path, libraries: dict[str, Library], forced: str | None) -> str:
    if forced:
        if forced not in libraries:
            raise IngestSkip(f"unknown_library:{forced}")
        return forced
    parts = set(path.parts)
    for library_id in libraries:
        if library_id in parts:
            return library_id
    return "model_methods" if "model_methods" in libraries else next(iter(libraries))


def extract_title(text: str, path: Path) -> str:
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            title = stripped.lstrip("#").strip()
            if title:
                return title[:160]
    for line in text.splitlines():
        stripped = line.strip()
        if stripped:
            return stripped[:160]
    return path.stem


def normalize_tags(value: Any, defaults: Iterable[str]) -> list[str]:
    items: list[str] = []
    for tag in defaults:
        if str(tag).strip():
            items.append(str(tag).strip())
    if isinstance(value, str):
        candidates = [value]
    elif isinstance(value, list):
        candidates = value
    else:
        candidates = []
    for tag in candidates:
        tag_text = str(tag).strip()
        if tag_text and tag_text not in items:
            items.append(tag_text)
    return items


def chunk_text(text: str, max_chars: int, min_chars: int) -> list[str]:
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if len(text) <= max_chars:
        return [text] if len(text) >= min_chars else []
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current.strip())
                current = ""
            for start in range(0, len(paragraph), max_chars):
                piece = paragraph[start : start + max_chars].strip()
                if len(piece) >= min_chars:
                    chunks.append(piece)
            continue
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if len(current) >= min_chars:
                chunks.append(current.strip())
            current = paragraph
    if len(current) >= min_chars:
        chunks.append(current.strip())
    return chunks


def stable_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def relpath_for(path: Path) -> str:
    try:
        return path.resolve().relative_to(REPO_ROOT).as_posix()
    except ValueError:
        return path.resolve().as_posix()


def make_chunks(path: Path, libraries: dict[str, Library], forced_library: str | None, max_chars: int, min_chars: int) -> list[Chunk]:
    raw_text = read_source_text(path)
    meta, body = parse_front_matter(raw_text)
    library_id = str(meta.get("library") or infer_library_id(path, libraries, forced_library))
    if library_id not in libraries:
        raise IngestSkip(f"unknown_library:{library_id}")
    library = libraries[library_id]
    title = str(meta.get("title") or extract_title(body, path))
    text_chunks = chunk_text(body, max_chars=max_chars, min_chars=min_chars)
    if not text_chunks:
        raise IngestSkip("no_chunk_after_min_chars")
    tags = normalize_tags(meta.get("tags"), library.default_tags)
    license_name = str(meta.get("license") or library.default_license)
    year = str(meta["year"]) if meta.get("year") is not None else None
    contest = str(meta["contest"]) if meta.get("contest") is not None else None
    problem_id = str(meta["problem_id"]) if meta.get("problem_id") is not None else None
    stage = str(meta.get("stage") or library.stage)
    recommended_use = str(
        meta.get("recommended_use")
        or f"Use as local evidence during {stage}; compare against the current problem before adopting."
    )
    risk_warning = str(
        meta.get("risk_warning")
        or "RAG evidence is advisory only; do not copy claims, formulas, or templates without data fit and source checks."
    )
    content_hash = stable_hash(raw_text)
    source_relpath = relpath_for(path)
    chunks: list[Chunk] = []
    for idx, chunk in enumerate(text_chunks):
        chunk_hash = stable_hash(chunk)
        chunk_id = stable_hash(f"{library_id}|{source_relpath}|{idx}|{chunk_hash}")[:32]
        chunks.append(
            Chunk(
                chunk_id=chunk_id,
                library_id=library_id,
                source_path=str(path.resolve()),
                source_relpath=source_relpath,
                content_hash=content_hash,
                chunk_hash=chunk_hash,
                chunk_index=idx,
                text=chunk,
                title=title,
                license=license_name,
                year=year,
                contest=contest,
                problem_id=problem_id,
                tags=tags,
                stage=stage,
                recommended_use=recommended_use,
                risk_warning=risk_warning,
            )
        )
    return chunks


def delete_existing_source(conn: sqlite3.Connection, library_id: str, source_relpath: str) -> None:
    existing = conn.execute(
        "SELECT chunk_id FROM chunks WHERE library_id = ? AND source_relpath = ?",
        (library_id, source_relpath),
    ).fetchall()
    if not existing:
        return
    ids = [row["chunk_id"] for row in existing]
    conn.executemany("DELETE FROM chunks WHERE chunk_id = ?", [(chunk_id,) for chunk_id in ids])
    try:
        conn.executemany("DELETE FROM chunks_fts WHERE chunk_id = ?", [(chunk_id,) for chunk_id in ids])
    except sqlite3.OperationalError:
        pass


def insert_chunks(conn: sqlite3.Connection, chunks: list[Chunk]) -> None:
    if not chunks:
        return
    delete_existing_source(conn, chunks[0].library_id, chunks[0].source_relpath)
    now = utc_now()
    rows = [
        (
            chunk.chunk_id,
            chunk.library_id,
            chunk.source_path,
            chunk.source_relpath,
            chunk.content_hash,
            chunk.chunk_hash,
            chunk.chunk_index,
            chunk.text,
            chunk.title,
            chunk.license,
            chunk.year,
            chunk.contest,
            chunk.problem_id,
            json.dumps(chunk.tags, ensure_ascii=False),
            chunk.stage,
            chunk.recommended_use,
            chunk.risk_warning,
            now,
        )
        for chunk in chunks
    ]
    conn.executemany(
        """
        INSERT OR REPLACE INTO chunks (
            chunk_id, library_id, source_path, source_relpath, content_hash,
            chunk_hash, chunk_index, text, title, license, year, contest,
            problem_id, tags_json, stage, recommended_use, risk_warning,
            ingested_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        rows,
    )
    try:
        conn.executemany(
            "INSERT INTO chunks_fts (chunk_id, title, text, tags) VALUES (?, ?, ?, ?)",
            [(chunk.chunk_id, chunk.title, chunk.text, " ".join(chunk.tags)) for chunk in chunks],
        )
    except sqlite3.OperationalError:
        pass


def log_error(conn: sqlite3.Connection, path: Path, library_id: str | None, reason: str, detail: str = "") -> None:
    conn.execute(
        """
        INSERT INTO ingest_errors (source_path, library_id, reason, detail, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (str(path), library_id, reason, detail[:1000], utc_now()),
    )


def hashing_embedding(text: str, dimensions: int = 384) -> list[float]:
    vector = [0.0] * dimensions
    tokens = re.findall(r"[A-Za-z0-9_]+|[\u4e00-\u9fff]", text.lower())
    for token in tokens:
        digest = hashlib.md5(token.encode("utf-8"), usedforsecurity=False).digest()
        index = int.from_bytes(digest[:4], "little") % dimensions
        vector[index] += 1.0
    norm = math.sqrt(sum(value * value for value in vector))
    if norm:
        vector = [value / norm for value in vector]
    return vector


def try_sentence_transformer_embeddings(texts: list[str], model_name: str, local_only: bool) -> list[list[float]] | None:
    try:
        from sentence_transformers import SentenceTransformer  # type: ignore
    except Exception:
        return None
    names = [model_name]
    if model_name != "BAAI/bge-small-zh-v1.5":
        names.append("BAAI/bge-small-zh-v1.5")
    for name in names:
        try:
            try:
                model = SentenceTransformer(name, local_files_only=local_only)
            except TypeError:
                if local_only:
                    return None
                model = SentenceTransformer(name)
            embeddings = model.encode(texts, normalize_embeddings=True)
            return [embedding.tolist() for embedding in embeddings]
        except Exception:
            continue
    return None


def update_chroma(chunks: list[Chunk], chroma_dir: Path, embedding_model: str, embedding_mode: str) -> str:
    if not chunks:
        return "skipped:no_chunks"
    try:
        import chromadb  # type: ignore
    except Exception:
        return "skipped:chromadb_not_installed"
    client = chromadb.PersistentClient(path=str(chroma_dir))
    collection = client.get_or_create_collection("mathmodel_knowledge")
    ids = [chunk.chunk_id for chunk in chunks]
    texts = [chunk.text for chunk in chunks]
    embeddings: list[list[float]] | None = None
    status_mode = "local-hashing-fallback"
    if embedding_mode in {"auto-local", "sentence-transformer"}:
        embeddings = try_sentence_transformer_embeddings(
            texts,
            embedding_model,
            local_only=embedding_mode == "auto-local",
        )
        if embeddings is not None:
            status_mode = embedding_model
    if embeddings is None:
        embeddings = [hashing_embedding(text) for text in texts]
        if embedding_mode == "sentence-transformer":
            status_mode = "local-hashing-fallback-after-transformer-error"
    try:
        collection.delete(ids=ids)
    except Exception:
        pass
    collection.add(
        ids=ids,
        documents=texts,
        embeddings=embeddings,
        metadatas=[
            {
                "library_id": chunk.library_id,
                "source_path": chunk.source_path,
                "source_relpath": chunk.source_relpath,
                "title": chunk.title,
                "license": chunk.license,
                "year": chunk.year or "",
                "contest": chunk.contest or "",
                "problem_id": chunk.problem_id or "",
                "tags": ",".join(chunk.tags),
                "stage": chunk.stage,
                "embedding_mode": status_mode,
            }
            for chunk in chunks
        ],
    )
    return f"updated:{status_mode}"


def ingest(args: argparse.Namespace) -> dict[str, Any]:
    source_paths = args.source or [REPO_ROOT / "knowledge" / "samples"]
    libraries = load_libraries(args.config)
    conn = init_db(args.db)
    summary: dict[str, Any] = {
        "db": str(args.db.resolve()),
        "sources": [str(path) for path in source_paths],
        "files_seen": 0,
        "files_ingested": 0,
        "chunks_written": 0,
        "files_skipped": 0,
        "errors": [],
        "chroma": "disabled" if args.vector_store == "none" else "not_requested",
    }
    chroma_chunks: list[Chunk] = []
    with conn:
        for path in iter_source_files(source_paths):
            summary["files_seen"] += 1
            try:
                chunks = make_chunks(
                    path,
                    libraries=libraries,
                    forced_library=args.library,
                    max_chars=args.max_chars,
                    min_chars=args.min_chars,
                )
                insert_chunks(conn, chunks)
                chroma_chunks.extend(chunks)
                summary["files_ingested"] += 1
                summary["chunks_written"] += len(chunks)
            except IngestSkip as exc:
                summary["files_skipped"] += 1
                reason = str(exc)
                log_error(conn, path, args.library, reason)
                summary["errors"].append({"path": str(path), "reason": reason})
            except Exception as exc:  # pragma: no cover - defensive logging
                summary["files_skipped"] += 1
                log_error(conn, path, args.library, "unexpected_error", repr(exc))
                summary["errors"].append({"path": str(path), "reason": "unexpected_error", "detail": repr(exc)})
    if args.vector_store in {"auto", "chroma"}:
        status = update_chroma(chroma_chunks, args.chroma_dir, args.embedding_model, args.embedding_mode)
        if args.vector_store == "chroma" and status.startswith("skipped"):
            raise SystemExit(f"Chroma indexing requested but {status}")
        summary["chroma"] = status
    return summary


def main() -> None:
    configure_stdout()
    args = parse_args()
    summary = ingest(args)
    if args.json:
        print(json.dumps(summary, ensure_ascii=False, indent=2))
    else:
        print(f"RAG ingest complete: {summary['files_ingested']} files, {summary['chunks_written']} chunks")
        print(f"SQLite ledger: {summary['db']}")
        if summary["files_skipped"]:
            print(f"Skipped: {summary['files_skipped']} files")
        if summary["chroma"] != "disabled":
            print(f"Chroma: {summary['chroma']}")


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        sys.exit(130)
